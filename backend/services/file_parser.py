from pathlib import Path

import fitz
from docx import Document
from pptx import Presentation


SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt"}


class FileParseError(ValueError):
    pass


def extract_text(path: Path) -> str:
    extension = path.suffix.lower()
    if extension not in SUPPORTED_EXTENSIONS:
        raise FileParseError(f"Unsupported file type: {extension}")

    if extension == ".txt":
        text = path.read_text(encoding="utf-8", errors="ignore")
    elif extension == ".pdf":
        text = _extract_pdf(path)
    elif extension == ".docx":
        text = _extract_docx(path)
    elif extension == ".pptx":
        text = _extract_pptx(path)
    else:
        raise FileParseError(f"Unsupported file type: {extension}")

    cleaned = "\n".join(line.strip() for line in text.splitlines() if line.strip())
    if not cleaned:
        raise FileParseError("No readable text could be extracted from this file.")
    return cleaned


def _extract_pdf(path: Path) -> str:
    chunks: list[str] = []
    with fitz.open(path) as document:
        for page in document:
            chunks.append(page.get_text())
    return "\n".join(chunks)


def _extract_docx(path: Path) -> str:
    document = Document(path)
    return "\n".join(paragraph.text for paragraph in document.paragraphs)


def _extract_pptx(path: Path) -> str:
    presentation = Presentation(path)
    chunks: list[str] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_text: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_text.append(shape.text)
        if slide_text:
            chunks.append(f"Slide {slide_number}\n" + "\n".join(slide_text))
    return "\n\n".join(chunks)
